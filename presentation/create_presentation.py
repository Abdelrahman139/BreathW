"""
BreathWise Graduation Presentation Generator
Generates a professional 21-slide PPTX and exports to PDF via PowerPoint COM.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Color Palette ──────────────────────────────────────────────
NAVY      = RGBColor(10, 25, 47)
DARK_NAVY = RGBColor(6, 15, 30)
WHITE     = RGBColor(255, 255, 255)
LIGHT     = RGBColor(200, 210, 220)
GRAY      = RGBColor(156, 163, 175)
BLUE      = RGBColor(56, 189, 248)
GREEN     = RGBColor(52, 211, 153)
RED       = RGBColor(248, 113, 113)
YELLOW    = RGBColor(251, 191, 36)
ORANGE    = RGBColor(251, 146, 60)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
IMG_DIR = os.path.join(os.path.dirname(__file__), "images")


def _img(name):
    """Return absolute path to an image in the images/ folder, or None."""
    p = os.path.join(IMG_DIR, name)
    return p if os.path.exists(p) else None


def _set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text,
              font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
              font_name="Arial"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_para(tf, text, font_size=16, color=WHITE, bold=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.space_before = space_before
    return p


def _add_bullet(tf, text, font_size=16, color=LIGHT, indent_level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.level = indent_level
    p.space_before = Pt(3)
    return p


def _divider(slide, top_inches, color=BLUE, width_inches=11, left_inches=1.2):
    from pptx.util import Inches as In
    shape = slide.shapes.add_shape(
        1,  # rectangle
        In(left_inches), In(top_inches),
        In(width_inches), In(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ══════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title Slide with landing page image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, DARK_NAVY)

    img = _img("landing_page.png")
    if img:
        slide.shapes.add_picture(img, Inches(6.5), Inches(0.5), width=Inches(6.5))

    tf = _add_text(slide, Inches(0.8), Inches(0.8), Inches(5.5), Inches(1),
                   "BreathWise", font_size=54, color=BLUE, bold=True)
    _add_para(tf, "AI Chest X-Ray Analysis and\nClinical Decision Support System",
              font_size=22, color=WHITE, bold=True)
    _add_para(tf, "", font_size=10)
    _add_para(tf, "University of Sadat City", font_size=16, color=LIGHT, bold=True)
    _add_para(tf, "Faculty of Computers and Artificial Intelligence", font_size=14, color=GRAY)
    _add_para(tf, "Department of Information Systems — Academic Year 2026", font_size=13, color=GRAY)
    _add_para(tf, "", font_size=8)
    _add_para(tf, "AbdAlrahman Magdy    |    Mahmoud Gamal    |    Seif Mohamed", font_size=12, color=LIGHT)
    _add_para(tf, "Mohamed Karam    |    Ahmed Eldamasy    |    Mahmoud Fadl", font_size=12, color=LIGHT)
    _add_para(tf, "", font_size=8)
    _add_para(tf, "Supervised by:", font_size=13, color=BLUE, bold=True)
    _add_para(tf, "Prof. Engy Elshafey   &   Prof. Asmaa Saad", font_size=13, color=LIGHT)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    tf = _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                   "The Problem", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    cards = [
        ("🏥  Radiologist Shortage & Burnout",
         "Imaging volume growing faster than radiologist training. Cognitive fatigue → diagnostic errors."),
        ("⏱️  Diagnostic Delay in Emergencies",
         "Hours of delay detecting Pneumothorax in ER/ICU can be fatal."),
        ("🔀  Inter-Observer Variability",
         "Same X-ray, different diagnoses depending on the reader and experience level."),
        ("🔒  The 'Black Box' Barrier",
         "Clinicians distrust AI that can't visually explain its reasoning."),
    ]
    for i, (title, desc) in enumerate(cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 6)
        top = Inches(1.8 + row * 2.6)
        tf2 = _add_text(slide, left, top, Inches(5.5), Inches(2),
                        title, font_size=20, color=WHITE, bold=True)
        _add_para(tf2, desc, font_size=14, color=GRAY)


def slide_03_disease_burden(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Disease Burden — Why These 4 Conditions?", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.2)

    diseases = [
        ("Pneumonia", "450M", "Leading infectious killer in children under 5 (740K deaths/yr)"),
        ("Cardiomegaly", "64M", "Heart failure, cardiomyopathy — high mortality & hospitalization"),
        ("Pleural Effusion", "10M", "Secondary to heart failure, pneumonia, and malignancies"),
        ("Pneumothorax", "0.5M", "Acute life-threatening emergency requiring immediate intervention"),
    ]
    for i, (name, cases, desc) in enumerate(diseases):
        top = Inches(1.7 + i * 1.35)
        tf = _add_text(slide, Inches(1), top, Inches(3), Inches(1),
                       name, font_size=22, color=WHITE, bold=True)
        _add_text(slide, Inches(4.5), top, Inches(2), Inches(1),
                  f"{cases} cases/yr", font_size=22, color=YELLOW, bold=True)
        _add_text(slide, Inches(7), top, Inches(5.5), Inches(1),
                  desc, font_size=14, color=GRAY)


def slide_04_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Project Objectives", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    objs = [
        ("🎯  High-Accuracy Classification", "Multi-label prediction with >90% accuracy across all thoracic conditions"),
        ("⚡  Near Real-Time Inference", "<0.5s per scan (single pass), ~2s with Test-Time Augmentation"),
        ("👁️  Explainability by Design", "Grad-CAM heatmaps overlaying every prediction to break the 'black box'"),
        ("📋  Clinical Workspace", "Patient tracking, scan histories, automated PDF reports with QuestPDF"),
    ]
    for i, (title, desc) in enumerate(objs):
        top = Inches(1.7 + i * 1.3)
        tf = _add_text(slide, Inches(1), top, Inches(11), Inches(1),
                       title, font_size=22, color=WHITE, bold=True)
        _add_para(tf, desc, font_size=15, color=GRAY)


def slide_05_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "System Architecture (3-Tier Microservices)", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.2)

    tiers = [
        ("React Frontend", "Port 5173", "TypeScript, Vite, TailwindCSS\nRecharts, Axios", BLUE),
        ("ASP.NET Core 8\nBackend", "Port 5169", "C#, EF Core, SQL Server\nJWT + RBAC, QuestPDF", GREEN),
        ("FastAPI / PyTorch\nAI Service", "Port 8000", "Python, DenseNet-121\nGrad-CAM, OpenCV", YELLOW),
    ]
    for i, (name, port, tech, color) in enumerate(tiers):
        left = Inches(0.8 + i * 4.2)
        top = Inches(2)
        tf = _add_text(slide, left, top, Inches(3.8), Inches(4),
                       name, font_size=22, color=color, bold=True)
        _add_para(tf, port, font_size=14, color=GRAY)
        _add_para(tf, "", font_size=6)
        _add_para(tf, tech, font_size=14, color=LIGHT)

    _add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
              "Backend acts as secure proxy → AI service never exposed to public internet. Internal auth: Authorization: internal-key",
              font_size=13, color=GRAY)


def slide_06_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Technology Stack", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    cols = [
        ("Frontend", BLUE, ["React 19", "TypeScript", "Vite", "TailwindCSS", "Recharts", "Axios"]),
        ("Backend", GREEN, ["C# .NET 8", "ASP.NET Core 8", "Entity Framework Core", "SQL Server", "QuestPDF", "JWT + RBAC"]),
        ("AI Service", YELLOW, ["Python 3.12", "FastAPI + Uvicorn", "PyTorch (DenseNet-121)", "OpenCV + PIL", "torchvision", "Grad-CAM"]),
    ]
    for i, (title, color, items) in enumerate(cols):
        left = Inches(0.8 + i * 4.2)
        tf = _add_text(slide, left, Inches(1.6), Inches(3.8), Inches(5),
                       title, font_size=24, color=color, bold=True)
        for item in items:
            _add_bullet(tf, f"•  {item}", font_size=15, color=LIGHT)


def slide_07_data_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "End-to-End Data Flow", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    steps = [
        ("1", "Upload", "Doctor uploads X-ray\nvia React drag-and-drop"),
        ("2", "Save", "C# saves image as\nwwwroot/uploads/[GUID].png"),
        ("3", "AI Call", "C# sends internal POST\nto localhost:8000/predict"),
        ("4", "Inference", "Python runs DenseNet-121\n+ TTA (3 passes) + Grad-CAM"),
        ("5", "DB Write", "C# decodes base64 heatmap\nwrites ScanResult to SQL"),
        ("6", "Display", "React renders diagnosis\n+ interactive heatmap toggle"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.5 + i * 2.1)
        tf = _add_text(slide, left, Inches(1.8), Inches(1.9), Inches(4),
                       num, font_size=36, color=BLUE, bold=True)
        _add_para(tf, title, font_size=18, color=WHITE, bold=True)
        _add_para(tf, desc, font_size=12, color=GRAY)


def slide_08_densenet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "The AI Model — DenseNet-121", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    tf = _add_text(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(5),
                   "Dense Connectivity", font_size=22, color=GREEN, bold=True)
    _add_bullet(tf, "Each layer receives feature maps from ALL preceding layers", color=LIGHT)
    _add_bullet(tf, "Growth Rate k = 32 (new feature maps per layer)", color=LIGHT)
    _add_bullet(tf, "Transition Layers: 1×1 conv + 2×2 avg pool between blocks", color=LIGHT)
    _add_para(tf, "", font_size=6)
    _add_para(tf, "Classifier Head", font_size=22, color=GREEN, bold=True)
    _add_bullet(tf, "nn.Linear(1024, 5) — 5 output logits", color=LIGHT)
    _add_bullet(tf, "Loss: BCEWithLogitsLoss (Binary Cross-Entropy)", color=LIGHT)
    _add_bullet(tf, "Optimizer: Adam, LR = 1e-4, Batch Size = 32", color=LIGHT)
    _add_para(tf, "", font_size=6)
    _add_para(tf, "Why DenseNet over ResNet?", font_size=20, color=YELLOW, bold=True)
    _add_bullet(tf, "1024 output features (vs ResNet's 2048) → more compact", color=LIGHT)
    _add_bullet(tf, "Better gradient flow through all layers", color=LIGHT)


def slide_09_gradcam(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
              "Explainability — Grad-CAM", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    tf = _add_text(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(2.5),
                   "Breaking the Black Box", font_size=20, color=GREEN, bold=True)
    _add_bullet(tf, "Computes gradients of class score vs DenseBlock4 activations", color=LIGHT, font_size=13)
    _add_bullet(tf, "Global Average Pooling → importance weights α", color=LIGHT, font_size=13)
    _add_bullet(tf, "Weighted sum → ReLU → upsampled to 224×224", color=LIGHT, font_size=13)
    _add_bullet(tf, "Mapped to JET color palette, blended α=0.4", color=LIGHT, font_size=13)
    _add_bullet(tf, "3 View Modes: Original | Overlay | Side-by-Side", color=LIGHT, font_size=13)

    img = _img("analysis_results.png")
    if img:
        slide.shapes.add_picture(img, Inches(6), Inches(1.2), width=Inches(7))


def slide_10_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Dataset & Class Balancing", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    tf = _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                   "Original NIH ChestX-ray14", font_size=22, color=YELLOW, bold=True)
    _add_bullet(tf, "112,120 images — 30,805 patients", color=LIGHT)
    _add_bullet(tf, "15 finding labels (14 pathologies + No Finding)", color=LIGHT)
    _add_bullet(tf, "Heavily imbalanced: No Finding = 60,361 (53.8%)", color=LIGHT)
    _add_bullet(tf, "Rarest: Hernia = 227 (0.2%)", color=LIGHT)
    _add_para(tf, "", font_size=8)
    _add_para(tf, "Problem: Training on raw data → biased model that always predicts 'Healthy'",
              font_size=14, color=RED)

    tf2 = _add_text(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5),
                    "Filtered & Balanced (BreathWise)", font_size=22, color=GREEN, bold=True)
    _add_bullet(tf2, "Filtered to 5 target conditions", color=LIGHT)
    _add_bullet(tf2, "Balanced to 5,500 per class", color=LIGHT)
    _add_bullet(tf2, "Total: 27,500 images", color=LIGHT)
    _add_bullet(tf2, "Train: 22,000 (80%) / Val: 5,500 (20%)", color=LIGHT)
    _add_para(tf2, "", font_size=8)
    _add_para(tf2, "Pneumonia (1,431 → 5,500) via oversampling with replacement",
              font_size=13, color=GRAY)


def slide_11_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
              "Training & Results", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    img = _img("training_results.png")
    if img:
        slide.shapes.add_picture(img, Inches(0.5), Inches(1.2), width=Inches(12.3))


def slide_12_per_class(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Per-Class Performance (Classification Report)", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.2)

    headers = ["Condition", "Precision", "Recall", "F1-Score", "Support"]
    rows = [
        ("Cardiomegaly ★", "0.82", "0.85", "0.84", "1,250"),
        ("Pneumonia", "0.85", "0.77", "0.81", "1,148"),
        ("Pneumothorax", "0.79", "0.76", "0.77", "1,194"),
        ("Pleural Effusion", "0.79", "0.72", "0.75", "1,928"),
        ("No Finding", "0.73", "0.66", "0.69", "1,113"),
        ("Macro Avg", "0.80", "0.75", "0.77", "6,633"),
    ]
    tf = _add_text(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(0.6),
                   "    ".join(f"{h:<16}" for h in headers),
                   font_size=16, color=BLUE, bold=True)
    for row in rows:
        color = GREEN if "★" in row[0] or row[0] == "Macro Avg" else LIGHT
        _add_para(tf, "    ".join(f"{c:<16}" for c in row), font_size=15, color=color)

    _add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1),
              "★ Cardiomegaly = highest F1 (0.84) due to distinct visual signature of enlarged heart\n"
              "  No Finding = lowest F1 (0.69) — expected, as it represents absence of features. Mitigated by mathematical derivation in production.",
              font_size=13, color=GRAY)


def slide_13_optimizations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Post-Training Optimizations", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    # Card 1: Thresholds
    tf = _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                   "Optimal Threshold Tuning", font_size=22, color=GREEN, bold=True)
    _add_bullet(tf, "Default 0.5 → sweep 0.0–1.0, maximize F1 per class", color=LIGHT)
    _add_para(tf, "", font_size=4)
    thresholds = [("Pneumonia", "0.39"), ("Effusion", "0.45"), ("Cardiomegaly", "0.54"),
                  ("Pneumothorax", "0.39"), ("No Finding", "0.43")]
    for name, val in thresholds:
        _add_bullet(tf, f"{name}: {val}", color=LIGHT, font_size=14)
    _add_para(tf, "", font_size=4)
    _add_para(tf, "Acute conditions get lower thresholds → higher sensitivity → fewer missed diagnoses",
              font_size=12, color=YELLOW)

    # Card 2: TTA
    tf2 = _add_text(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5),
                    "Test-Time Augmentation (TTA)", font_size=22, color=GREEN, bold=True)
    _add_bullet(tf2, "3 predictions per image:", color=LIGHT)
    _add_bullet(tf2, "Original + H-Flip + 10° Rotation", color=LIGHT, indent_level=1)
    _add_bullet(tf2, "Final = arithmetic mean of 3 sigmoid outputs", color=LIGHT)
    _add_para(tf2, "", font_size=8)
    _add_para(tf2, "Progressive Accuracy:", font_size=18, color=YELLOW, bold=True)
    _add_bullet(tf2, "Baseline (0.5): 89.37%", color=LIGHT, font_size=14)
    _add_bullet(tf2, "+ Thresholds: 89.33%", color=LIGHT, font_size=14)
    _add_bullet(tf2, "+ TTA: 89.67% ✓", color=GREEN, font_size=14)


def slide_14_dashboard(prs):
    """Dedicated Dashboard slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
              "Doctor Dashboard — Clinical Workspace", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    img = _img("dashboard.png")
    if img:
        slide.shapes.add_picture(img, Inches(1), Inches(1.3), width=Inches(11.3))

    _add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.6),
              "KPI Cards (Total Patients, Total Scans, High Risk Cases) • Detection Distribution Chart • Healthy vs At Risk Donut • Recent Scans with Risk Badges • Scans Over Time (Last 7 Days)",
              font_size=12, color=GRAY)


def slide_15_analysis_results(prs):
    """Dedicated Analysis Results slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
              "AI Diagnosis & Grad-CAM Results", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    img = _img("analysis_results.png")
    if img:
        slide.shapes.add_picture(img, Inches(1), Inches(1.3), width=Inches(11.3))

    _add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.6),
              "Grad-CAM Side-by-Side View • Original X-ray vs Heatmap Overlay • AI Diagnosis Panel with Risk Badges • Download PDF Report Button • Doctor's Notes Section",
              font_size=12, color=GRAY)


def slide_16_pdf_report(prs):
    """PDF Report slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
              "Automated Clinical PDF Reports", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    tf = _add_text(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(5),
                   "QuestPDF Report Engine", font_size=20, color=GREEN, bold=True)
    _add_bullet(tf, "Professional radiology report layout", color=LIGHT)
    _add_bullet(tf, "Patient demographics & study info", color=LIGHT)
    _add_bullet(tf, "AI Predictive Analysis with probability bars", color=LIGHT)
    _add_bullet(tf, "Overall Risk Severity Gauge (Low / Moderate / High)", color=LIGHT)
    _add_bullet(tf, "Probability Distribution Chart", color=LIGHT)
    _add_bullet(tf, "Clinical Impression with recommendations", color=LIGHT)
    _add_bullet(tf, "Attached X-ray + Grad-CAM heatmap images", color=LIGHT)
    _add_para(tf, "", font_size=6)
    _add_para(tf, "Generated via C# QuestPDF — instant compilation", font_size=13, color=GRAY)
    _add_para(tf, "Requires clinician's electronic signature (human-in-the-loop)", font_size=13, color=YELLOW)

    img = _img("pdf_report.png")
    if img:
        slide.shapes.add_picture(img, Inches(6.5), Inches(1.0), height=Inches(6))


def slide_17_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Security & Regulatory Compliance", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    pillars = [
        ("🔐  JWT Authentication", "24-hour token expiry, validated on every API request"),
        ("👥  Role-Based Access (RBAC)", "Doctor (full CRUD), Patient (read-only own data), Admin"),
        ("🔑  PBKDF2 Password Hashing", "HMAC-SHA256 + unique 128-bit salt per user, 10,000 iterations"),
        ("🏷️  Anonymized GUID Filenames", "X-rays saved as 8c4e2a87-19db-4382.png — no patient names on disk"),
        ("📜  HIPAA/GDPR Alignment", "Query-level isolation matches JWT identity claims. 17 secured API endpoints"),
    ]
    for i, (title, desc) in enumerate(pillars):
        top = Inches(1.6 + i * 1.1)
        tf = _add_text(slide, Inches(1), top, Inches(11), Inches(0.9),
                       title, font_size=20, color=WHITE, bold=True)
        _add_para(tf, desc, font_size=14, color=GRAY)


def slide_18_graceful(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Reliability — Graceful Degradation", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    steps = [
        ("1", "AI Service DOWN", "Network timeout, crash,\nor GPU unavailable", RED),
        ("2", "C# Catch Block", "AiService catches exception\n→ falls back to GetMockScores()", YELLOW),
        ("3", "UI Uninterrupted", "Doctors can still log in,\nview patients, browse dashboard", GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        left = Inches(1 + i * 4)
        tf = _add_text(slide, left, Inches(2), Inches(3.5), Inches(3),
                       num, font_size=48, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        _add_para(tf, title, font_size=22, color=WHITE, bold=True)
        _add_para(tf, desc, font_size=14, color=GRAY)

    _add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.5),
              "Also: Python AI service has its own demo mode — if densenet121.pt weights are missing, it serves mock predictions so frontend testing can proceed.",
              font_size=14, color=GRAY)


def slide_19_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
              "Limitations & Future Roadmap", font_size=40, color=BLUE, bold=True)
    _divider(slide, 1.2)

    # Limitations
    tf = _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5),
                   "Current Limitations", font_size=22, color=RED, bold=True)
    lims = [
        "Dataset scope — NIH only, may vary on different equipment",
        "JPEG/PNG input only — no native DICOM support",
        "Single modality — chest X-rays only (no CT/MRI)",
        "No real-time multi-physician collaboration",
        "On-premise deployment only",
    ]
    for l in lims:
        _add_bullet(tf, f"✗  {l}", color=LIGHT, font_size=13)

    # Roadmap
    tf2 = _add_text(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5),
                    "Future Roadmap", font_size=22, color=GREEN, bold=True)
    roads = [
        "DICOM & PACS Integration (fo-dicom C-STORE)",
        "VLM Report Writing (LLaVA-Med / Med-PaLM)",
        "More Pathologies: TB (1.3M deaths/yr), Atelectasis, COPD",
        "Federated Learning across hospitals",
        "Cloud-Native auto-scaling deployment",
    ]
    for r in roads:
        _add_bullet(tf2, f"→  {r}", color=LIGHT, font_size=13)


def slide_20_diseases_page(prs):
    """Diseases education page slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
              "Patient Education — Disease Knowledge Base", font_size=36, color=BLUE, bold=True)
    _divider(slide, 1.0)

    img = _img("diseases_page.png")
    if img:
        slide.shapes.add_picture(img, Inches(1), Inches(1.3), width=Inches(11.3))

    _add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.6),
              "Comprehensive disease information • Common symptoms & causes • When to See a Doctor alerts • Patient-accessible (read-only role)",
              font_size=12, color=GRAY)


def slide_21_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_NAVY)

    tf = _add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                   "BreathWise", font_size=54, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_para(tf, "AI won't replace radiologists, but radiologists\nwho use AI will replace those who don't.",
              font_size=20, color=LIGHT)
    _add_para(tf, "", font_size=10)

    objectives = [
        "✅ Multi-label classification (AUC 0.91–0.96)",
        "✅ Explainable AI — Grad-CAM in production",
        "✅ Secure RBAC on all 17 API endpoints",
        "✅ Automated PDF reports (QuestPDF)",
        "✅ Real-time dashboard (7 live metrics)",
        "✅ Graceful degradation & demo mode",
        "✅ HIPAA-aligned privacy (GUID, PBKDF2, JWT)",
        "✅ Modular, extensible REST architecture",
    ]
    for obj in objectives:
        _add_para(tf, obj, font_size=14, color=GREEN)

    _add_para(tf, "", font_size=6)
    _add_para(tf, "Thank You!", font_size=36, color=WHITE, bold=True)
    _add_para(tf, "Prof. Engy Elshafey  &  Prof. Asmaa Saad", font_size=16, color=GRAY)


# ══════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)          # 1  Title
    slide_02_problem(prs)        # 2  The Problem
    slide_03_disease_burden(prs) # 3  Disease Burden
    slide_04_objectives(prs)     # 4  Project Objectives
    slide_05_architecture(prs)   # 5  Architecture
    slide_06_tech_stack(prs)     # 6  Tech Stack
    slide_07_data_flow(prs)      # 7  Data Flow
    slide_08_densenet(prs)       # 8  DenseNet-121
    slide_09_gradcam(prs)        # 9  Grad-CAM
    slide_10_dataset(prs)        # 10 Dataset
    slide_11_training(prs)       # 11 Training (screenshot)
    slide_12_per_class(prs)      # 12 Per-Class Metrics
    slide_13_optimizations(prs)  # 13 Optimizations
    slide_14_dashboard(prs)      # 14 Dashboard (screenshot)
    slide_15_analysis_results(prs)  # 15 Analysis Results (screenshot)
    slide_16_pdf_report(prs)     # 16 PDF Report (screenshot)
    slide_17_security(prs)       # 17 Security
    slide_18_graceful(prs)       # 18 Graceful Degradation
    slide_19_limitations(prs)    # 19 Limitations & Roadmap
    slide_20_diseases_page(prs)  # 20 Disease Education (screenshot)
    slide_21_thankyou(prs)       # 21 Thank You

    out_pptx = "BreathWise_Presentation.pptx"
    prs.save(out_pptx)
    print(f"[OK] PPTX saved: {out_pptx}  ({len(prs.slides)} slides)")

    # Export to PDF via PowerPoint COM
    try:
        import comtypes.client
        abs_pptx = os.path.abspath(out_pptx)
        abs_pdf  = os.path.abspath(out_pptx.replace(".pptx", ".pdf"))
        print("Opening PowerPoint via COM...")
        pp = comtypes.client.CreateObject("Powerpoint.Application")
        pp.Visible = 1
        deck = pp.Presentations.Open(abs_pptx)
        deck.SaveAs(abs_pdf, 32)
        deck.Close()
        pp.Quit()
        print(f"[OK] PDF saved: {abs_pdf}")
    except Exception as e:
        print(f"[WARN] PDF export skipped (PowerPoint not available): {e}")


if __name__ == "__main__":
    main()
